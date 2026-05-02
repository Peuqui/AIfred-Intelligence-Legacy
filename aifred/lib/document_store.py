"""
Document Store - Chunking, Embedding & ChromaDB Storage for uploaded documents.

Uses the same ChromaDB server and Ollama embedding function as vector_cache.py
but with a separate collection for user documents.
"""

from datetime import datetime
from pathlib import Path
from typing import Any, Optional

import chromadb
from chromadb.config import Settings
from chromadb.errors import ChromaError

from .config import (
    DEFAULT_OLLAMA_URL,
    DOCUMENT_CHUNK_OVERLAP,
    DOCUMENT_CHUNK_SIZE,
    DOCUMENT_COLLECTION,
    DOCUMENTS_DIR,
)
from .logging_utils import log_message
from .vector_cache import OLLAMA_EMBEDDING_MODEL, OllamaEmbeddingFunction


def _read_text_file(file_path: Path) -> str:
    """Read a text file with automatic encoding detection via chardet."""
    raw = file_path.read_bytes()
    import chardet
    detected = chardet.detect(raw)
    encoding = detected.get("encoding") or "utf-8"
    return raw.decode(encoding)


def _read_pdf(file_path: Path) -> str:
    """Extract text from a PDF using pdftotext (poppler-utils).

    pdftotext joins hyphenated line-breaks back into whole words and
    converts ligatures (ﬂ, ﬁ) to plain letters — both crucial for
    embedding quality. PyMuPDF/fitz preserves them as-is which results
    in fragmented embeddings (`Misch-` + `volk` as two halves, `Brotﬂ aden`
    as a strange unicode word).
    """
    import shutil
    import subprocess
    if not shutil.which("pdftotext"):
        raise RuntimeError(
            "pdftotext not installed. Please install poppler-utils:\n"
            "  Debian/Ubuntu:  sudo apt install poppler-utils\n"
            "  Fedora/RHEL:    sudo dnf install poppler-utils\n"
            "  macOS:          brew install poppler"
        )
    result = subprocess.run(
        ["pdftotext", "-enc", "UTF-8", str(file_path), "-"],
        capture_output=True,
        text=True,
        check=True,
        timeout=300,
    )
    return result.stdout


def _read_csv(file_path: Path) -> str:
    """Read a CSV file and return as markdown table."""
    import csv
    import io
    content = _read_text_file(file_path)
    reader = csv.reader(io.StringIO(content))
    rows = list(reader)
    if not rows:
        return ""
    # Header + separator + data rows
    header = "| " + " | ".join(rows[0]) + " |"
    separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
    data = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
    return f"{header}\n{separator}\n{data}"


def _read_docx(file_path: Path) -> str:
    """Extract text from a .docx file."""
    from docx import Document
    doc = Document(str(file_path))
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_xlsx(file_path: Path) -> str:
    """Extract text from a .xlsx file as markdown tables (one per sheet)."""
    from openpyxl import load_workbook
    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = [[str(cell) if cell is not None else "" for cell in row] for row in ws.iter_rows(values_only=True)]
        if not rows:
            continue
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
        data = "\n".join("| " + " | ".join(row) + " |" for row in rows[1:])
        parts.append(f"## {sheet}\n\n{header}\n{separator}\n{data}")
    wb.close()
    return "\n\n".join(parts)


def _read_pptx(file_path: Path) -> str:
    """Extract text from a .pptx file."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            parts.append(f"## Slide {i}\n\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _read_odt(file_path: Path) -> str:
    """Extract text from a .odt file."""
    from odf.opendocument import load
    from odf.text import P
    doc = load(str(file_path))
    paragraphs = doc.getElementsByType(P)
    return "\n\n".join(
        "".join(str(node) for node in p.childNodes)
        for p in paragraphs
        if p.childNodes
    )


def _read_ods(file_path: Path) -> str:
    """Extract text from a .ods spreadsheet as markdown tables."""
    from odf.opendocument import load
    from odf.table import Table, TableRow, TableCell
    from odf.text import P
    doc = load(str(file_path))
    parts: list[str] = []
    for table in doc.getElementsByType(Table):
        name = table.getAttribute("name") or "Sheet"
        rows: list[list[str]] = []
        for row in table.getElementsByType(TableRow):
            cells: list[str] = []
            for cell in row.getElementsByType(TableCell):
                text = "".join(
                    "".join(str(n) for n in p.childNodes)
                    for p in cell.getElementsByType(P)
                )
                cells.append(text)
            if any(c.strip() for c in cells):
                rows.append(cells)
        if not rows:
            continue
        # Normalize column count
        max_cols = max(len(r) for r in rows)
        rows = [r + [""] * (max_cols - len(r)) for r in rows]
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
        data = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
        parts.append(f"## {name}\n\n{header}\n{separator}\n{data}")
    return "\n\n".join(parts)


def _read_odp(file_path: Path) -> str:
    """Extract text from a .odp presentation."""
    from odf.opendocument import load
    from odf.draw import Page
    from odf.text import P
    doc = load(str(file_path))
    parts: list[str] = []
    for i, page in enumerate(doc.getElementsByType(Page), 1):
        texts: list[str] = []
        for p in page.getElementsByType(P):
            text = "".join(str(n) for n in p.childNodes)
            if text.strip():
                texts.append(text)
        if texts:
            parts.append(f"## Slide {i}\n\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Split text into chunks of *exactly* chunk_size tokens.

    Uses the Qwen3 tokenizer (already cached locally for token estimation
    elsewhere in the project) for true token-count accuracy. This replaces
    the old char-heuristic which under-estimated tokens for token-dense
    languages (German + Hebrew inserts) and could push embeddings past
    the model's context limit.

    Falls back to char-based slicing only if the tokenizer cannot be
    loaded — that path is safe but less accurate.
    """
    try:
        return _chunk_by_tokens(text, chunk_size, overlap)
    except Exception as exc:
        log_message(f"⚠️ Tokenizer chunking failed ({exc}) — using char fallback")
        return _chunk_by_chars(text, chunk_size, overlap)


def _chunk_by_tokens(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Token-accurate chunker using the Qwen3 tokenizer."""
    from .context_manager import _tokenizer_cache, count_tokens_with_tokenizer

    # Warm the cache via the public counter (handles download + caching)
    count_tokens_with_tokenizer("warmup")
    tokenizer = _tokenizer_cache.get("qwen3")
    if tokenizer is None:
        raise RuntimeError("Qwen3 tokenizer not initialised")

    encoding = tokenizer.encode(text)
    token_ids = encoding.ids
    if len(token_ids) <= chunk_size:
        stripped = text.strip()
        return [stripped] if stripped else []

    chunks: list[str] = []
    start = 0
    while start < len(token_ids):
        end = min(start + chunk_size, len(token_ids))
        piece = tokenizer.decode(token_ids[start:end]).strip()
        if piece:
            chunks.append(piece)
        if end >= len(token_ids):
            break
        start = end - overlap
    return chunks


def _chunk_by_chars(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Char-based fallback. Less accurate, used only when tokenizer unavailable."""
    from .config import CHARS_PER_TOKEN
    chunk_chars = chunk_size * CHARS_PER_TOKEN
    overlap_chars = overlap * CHARS_PER_TOKEN

    if len(text) <= chunk_chars:
        return [text]

    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + chunk_chars
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap_chars
    return chunks


PARSERS = {
    ".pdf": _read_pdf,
    ".txt": _read_text_file,
    ".md": _read_text_file,
    ".csv": _read_csv,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
    ".odt": _read_odt,
    ".ods": _read_ods,
    ".odp": _read_odp,
}


class DocumentStore:
    """Manages document chunking, embedding and retrieval via ChromaDB."""

    def __init__(self, host: str = "localhost", port: int = 8000):
        self._client = chromadb.HttpClient(
            host=host, port=port,
            settings=Settings(anonymized_telemetry=False),
        )
        self._client.heartbeat()
        # Two embedding functions for two access patterns:
        # - bulk indexing (many chunks at once) → GPU + 10 min keep_alive
        #   so the model stays warm for the duration of the upload
        # - single-shot query (one user question) → CPU, no GPU residency
        #   so we don't fight the active LLM for VRAM
        self._embed_index = OllamaEmbeddingFunction(
            model_name=OLLAMA_EMBEDDING_MODEL,
            host=DEFAULT_OLLAMA_URL,
            mode="index",
        )
        self._embed_query = OllamaEmbeddingFunction(
            model_name=OLLAMA_EMBEDDING_MODEL,
            host=DEFAULT_OLLAMA_URL,
            mode="query",
        )
        # Collection-default uses the query function — anything that doesn't
        # explicitly pass embeddings (only happens accidentally) stays on CPU.
        self._embed_fn = self._embed_query
        self._collection = self._client.get_or_create_collection(
            name=DOCUMENT_COLLECTION,
            metadata={
                "description": "AIfred uploaded documents",
                "embedding_model": OLLAMA_EMBEDDING_MODEL,
            },
            embedding_function=self._embed_fn,  # type: ignore[arg-type]
        )
        DOCUMENTS_DIR.mkdir(parents=True, exist_ok=True)
        log_message(f"📄 DocumentStore connected: {self._collection.count()} chunks")

    async def index_document(self, file_path: Path, filename: str) -> int:
        """Parse, chunk and embed a document. Returns number of chunks created."""
        suffix = file_path.suffix.lower()
        parser = PARSERS.get(suffix)
        if not parser:
            raise ValueError(f"Unsupported file type: {suffix}")

        text = parser(file_path)
        if not text.strip():
            raise ValueError(f"Document is empty: {filename}")

        chunks = _chunk_text(text, DOCUMENT_CHUNK_SIZE, DOCUMENT_CHUNK_OVERLAP)
        now = datetime.now().isoformat()

        # Folder is the parent path of the relative filename — empty string means root
        folder = str(Path(filename).parent) if "/" in filename else ""

        ids = [f"{filename}__chunk_{i}" for i in range(len(chunks))]
        metadatas = [
            {
                "filename": filename,
                "folder": folder,
                "chunk_index": i,
                "total_chunks": len(chunks),
                "upload_date": now,
                "source_trust": "trusted",  # Documents uploaded via UI are trusted
            }
            for i in range(len(chunks))
        ]

        # Delete any existing chunks for this filename first. Without this,
        # re-indexing a now-shorter file would leave zombie chunks in place
        # (old chunk_N where N > new total_chunks). With delete-before-upsert
        # the file always reflects the current source — no stale data.
        self._collection.delete(where={"filename": filename})

        # Embed via the index-mode function (GPU + warm cache). Passing
        # the embeddings explicitly bypasses the collection's default
        # embedding_function (which is the query/CPU one) for this write.
        embeddings = self._embed_index(chunks)
        self._collection.upsert(
            ids=ids,
            documents=chunks,
            embeddings=embeddings,
            metadatas=metadatas,  # type: ignore[arg-type]
        )

        log_message(f"📄 Indexed {filename}: {len(chunks)} chunks")
        return len(chunks)

    async def search(
        self,
        query: str,
        n_results: int = 5,
        folder: Optional[str] = None,
        neighbor_window: Optional[int] = None,
    ) -> list[dict[str, Any]]:
        """Semantic search across documents. Returns list of chunk results.

        Args:
            query: Search query.
            n_results: Max number of chunks to return.
            folder: If set, restrict search to this folder. Sub-folders are
                    NOT included — pass the exact folder string as stored
                    in metadata (e.g. "bibel/Schlachter").
            neighbor_window: Per hit, also include ±N adjacent chunks of the
                    same document so the model sees the full surrounding
                    context (mid-sentence chunk cuts are mitigated).
                    None → use DOCUMENT_SEARCH_NEIGHBOR_WINDOW from config,
                    0 → off (similarity-only).
        """
        if self._collection.count() == 0:
            return []

        if neighbor_window is None:
            from .config import DOCUMENT_SEARCH_NEIGHBOR_WINDOW
            neighbor_window = DOCUMENT_SEARCH_NEIGHBOR_WINDOW

        # Embed the query via the CPU/query-mode function so we don't
        # wake up the GPU embedding model for a single-shot search.
        query_embeddings = self._embed_query([query])
        query_kwargs: dict[str, Any] = {
            "query_embeddings": query_embeddings,
            "n_results": min(n_results, self._collection.count()),
        }
        if folder is not None:
            query_kwargs["where"] = {"folder": folder}

        results = self._collection.query(**query_kwargs)

        hits: list[dict[str, Any]] = []
        if results and results["documents"] and results["documents"][0]:
            for i, doc in enumerate(results["documents"][0]):
                meta = results["metadatas"][0][i] if results["metadatas"] else {}  # type: ignore[index]
                distance = results["distances"][0][i] if results["distances"] else None  # type: ignore[index]
                hits.append({
                    "content": doc,
                    "filename": meta.get("filename", ""),
                    "folder": meta.get("folder", ""),
                    "chunk_index": meta.get("chunk_index", 0),
                    "total_chunks": meta.get("total_chunks", 0),
                    "distance": distance,
                    "_neighbor": False,  # mark original similarity hits
                })

        # Fetch neighbor chunks for each similarity hit (deduped by ID).
        # Neighbors are added with _neighbor=True so callers can distinguish
        # similarity-driven hits from contextual augmentation.
        if neighbor_window and neighbor_window > 0 and hits:
            existing_ids = {
                f"{h['filename']}__chunk_{h['chunk_index']}" for h in hits
            }
            neighbor_ids: set[str] = set()
            for h in hits:
                fname = h["filename"]
                idx = h["chunk_index"]
                total = h.get("total_chunks", 0)
                for offset in range(-neighbor_window, neighbor_window + 1):
                    if offset == 0:
                        continue
                    nidx = idx + offset
                    if nidx < 0 or (total and nidx >= total):
                        continue
                    nid = f"{fname}__chunk_{nidx}"
                    if nid not in existing_ids:
                        neighbor_ids.add(nid)

            if neighbor_ids:
                neighbor_data = self._collection.get(ids=list(neighbor_ids))
                if neighbor_data and neighbor_data.get("documents"):
                    for i, doc in enumerate(neighbor_data["documents"]):
                        meta = (neighbor_data["metadatas"][i]
                                if neighbor_data.get("metadatas") else {})  # type: ignore[index]
                        hits.append({
                            "content": doc,
                            "filename": meta.get("filename", ""),
                            "folder": meta.get("folder", ""),
                            "chunk_index": meta.get("chunk_index", 0),
                            "total_chunks": meta.get("total_chunks", 0),
                            "distance": None,  # not from similarity search
                            "_neighbor": True,
                        })

            # Sort by (filename, chunk_index) so consecutive chunks are
            # adjacent in the output — easier for the model to reconstruct
            # the original passage.
            hits.sort(key=lambda h: (h["filename"], h["chunk_index"]))

        return hits

    def list_documents(self) -> list[dict[str, Any]]:
        """List all unique documents with metadata."""
        if self._collection.count() == 0:
            return []

        all_data = self._collection.get()
        docs: dict[str, dict[str, Any]] = {}
        if all_data and all_data["metadatas"]:
            for meta in all_data["metadatas"]:
                if not isinstance(meta, dict):
                    continue
                fname = str(meta.get("filename", ""))
                if fname and fname not in docs:
                    docs[fname] = {
                        "filename": fname,
                        "total_chunks": meta.get("total_chunks", 0),
                        "upload_date": meta.get("upload_date", ""),
                    }

        return list(docs.values())

    async def delete_document(self, filename: str, delete_file: bool = True) -> int:
        """Delete all chunks for a document. Returns number of deleted chunks."""
        all_data = self._collection.get(
            where={"filename": filename},
        )
        if not all_data or not all_data["ids"]:
            return 0

        count = len(all_data["ids"])
        self._collection.delete(ids=all_data["ids"])

        # Delete file from disk if requested
        if delete_file:
            file_path = DOCUMENTS_DIR / filename
            if file_path.exists():
                file_path.unlink()

        log_message(f"🗑️ Deleted {filename}: {count} chunks")
        return count


# Singleton
_store: Optional[DocumentStore] = None


def get_document_store() -> Optional[DocumentStore]:
    """Get or create the global DocumentStore singleton."""
    global _store
    if _store is None:
        try:
            _store = DocumentStore()
        except (ChromaError, ConnectionError, OSError, ValueError) as e:
            log_message(f"❌ DocumentStore init failed: {e}")
            return None
    return _store
